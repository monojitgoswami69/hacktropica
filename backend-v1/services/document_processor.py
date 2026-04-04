"""
Document Processing Service
PDF, DOCX, PPTX, Image extraction with OCR fallback.
"""

from __future__ import annotations

import io
import re
import logging
from typing import Dict, List, Optional, Tuple

from pypdf import PdfReader

logger = logging.getLogger(__name__)

# Optional OCR dependencies
try:
    import pytesseract
    from PIL import Image
    import pypdfium2 as pdfium

    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logger.warning("⚠️  OCR dependencies not available (pytesseract/Pillow/pypdfium2)")


class DocumentProcessor:
    """Stateless document processing — all methods are static/classmethod."""

    # ── PDF ────────────────────────────────────────────────────────────

    @staticmethod
    def extract_pdf_text(file_bytes: bytes) -> Tuple[str, List[Tuple[int, str]]]:
        reader = PdfReader(io.BytesIO(file_bytes))
        page_texts: List[str] = []
        page_data: List[Tuple[int, str]] = []

        for i, page in enumerate(reader.pages):
            try:
                text = page.extract_text() or ""
                if text.strip():
                    page_texts.append(text)
                    page_data.append((i + 1, text))
            except Exception:
                continue

        extracted = "\n\n".join(page_texts).strip()
        if not extracted:
            raise ValueError("No extractable text found in PDF")
        return extracted, page_data

    @staticmethod
    def extract_pdf_ocr(file_bytes: bytes) -> Tuple[str, List[Tuple[int, str]]]:
        if not OCR_AVAILABLE:
            raise RuntimeError("OCR dependencies not installed")

        doc = pdfium.PdfDocument(io.BytesIO(file_bytes))
        page_texts: List[str] = []
        page_data: List[Tuple[int, str]] = []

        for idx in range(len(doc)):
            try:
                page = doc[idx]
                bitmap = page.render(scale=2.0)
                pil_img = bitmap.to_pil().convert("RGB")
                text = pytesseract.image_to_string(pil_img)
                if text and text.strip():
                    page_texts.append(text.strip())
                    page_data.append((idx + 1, text.strip()))
            except Exception:
                continue

        extracted = "\n\n".join(page_texts).strip()
        if not extracted:
            raise ValueError("No extractable text found via OCR")
        return extracted, page_data

    # ── Image ─────────────────────────────────────────────────────────

    @staticmethod
    def extract_image_ocr(file_bytes: bytes) -> str:
        if not OCR_AVAILABLE:
            raise RuntimeError("OCR dependencies not installed")

        img = Image.open(io.BytesIO(file_bytes)).convert("RGB")
        text = pytesseract.image_to_string(img).strip()
        if not text:
            raise ValueError("No extractable text found in image")
        return text

    # ── DOCX ──────────────────────────────────────────────────────────

    @staticmethod
    def extract_docx_text(file_bytes: bytes) -> str:
        import docx

        doc = docx.Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        extracted = "\n\n".join(paragraphs)
        if not extracted:
            raise ValueError("DOCX contains no extractable text")
        return extracted

    # ── PPTX ──────────────────────────────────────────────────────────

    @staticmethod
    def extract_pptx_text(file_bytes: bytes) -> Tuple[str, List[Tuple[int, str]]]:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        page_data: List[Tuple[int, str]] = []

        for idx, slide in enumerate(prs.slides):
            parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if parts:
                page_data.append((idx + 1, "\n".join(parts)))

        if not page_data:
            raise ValueError("PPTX contains no extractable text")
        extracted = "\n\n".join(text for _, text in page_data)
        return extracted, page_data

    # ── Chunking ──────────────────────────────────────────────────────

    @staticmethod
    def chunk_text(
        text: str,
        chunk_size: int,
        overlap: int,
        page_texts: Optional[List[Tuple[int, str]]] = None,
    ) -> List[Dict[str, object]]:
        """Split text into chunks. Preserves page info when available."""

        def _chunk_string(input_text: str) -> List[str]:
            sentences = re.split(r"(?<=[.!?])\s+", input_text.strip())
            sentences = [s.strip() for s in sentences if s.strip()]
            if not sentences:
                return []

            out: List[str] = []
            current = ""
            for sentence in sentences:
                if len(current) + len(sentence) + 1 <= chunk_size:
                    current += (" " if current else "") + sentence
                else:
                    if current:
                        out.append(current)
                    current = sentence
            if current:
                out.append(current)
            return out

        if page_texts:
            all_chunks: List[Dict[str, object]] = []
            for page_num, page_text in page_texts:
                for chunk in _chunk_string(page_text):
                    all_chunks.append({"text": chunk, "page_start": page_num, "page_end": page_num})
            if all_chunks:
                return all_chunks

        return [{"text": chunk, "page_start": None, "page_end": None} for chunk in _chunk_string(text)]
